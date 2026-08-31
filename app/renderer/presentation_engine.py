"""
IAS Presentation Engine v3
──────────────────────────
Full pipeline: Document → intelligent topic grouping → rich slide planning → polished PPTX

Design principles (from PPTX skill):
  - Every slide has a visual element (never text-only)
  - Varied layouts — no two consecutive slides look the same
  - Per-theme color palettes: dominant/support/accent ("sandwich" dark→light→dark)
  - Icon circles for visual interest
  - Charts embedded as images where data supports them
  - NEVER: accent stripes, header bars, color bars down edges
  - Presenter mode: sparse text, large type, talking-point bullets
  - Detailed mode: full bullets, sub-items, context
"""

import re
import io
import math
from dataclasses import dataclass, field
from typing import List, Optional, Dict, Tuple
from app.models.document import Block, BlockType, Document, Section
from app.renderer.visual_generator import generate as gen_visual


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
#  THEME PALETTES
#  Per-theme: bg (light content), dark (title/conclusion),
#  primary, secondary, accent, text_light, text_dark, card
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

PALETTES: Dict[str, Dict] = {
    "academic": {
        "bg":       "F7F8FC", "dark":     "1A1A2E",
        "primary":  "2B4C9B", "secondary":"5D3A8E",
        "accent":   "C4952A", "muted":    "8892A4",
        "text_dark":"FFFFFF", "text_light":"1A1A2E",
        "card":     "EEF1F8", "card2":    "F2ECF8",
        "warn":     "C0392B", "success":  "1A6E4A",
        "font_h":   "Cambria", "font_b": "Calibri",
    },
    "magazine": {
        "bg":       "FFFFFF", "dark":     "111111",
        "primary":  "E63946", "secondary":"111111",
        "accent":   "F4A261", "muted":    "888888",
        "text_dark":"FFFFFF", "text_light":"111111",
        "card":     "FFF5F5", "card2":    "F5F5F5",
        "warn":     "C0392B", "success":  "2D6A4F",
        "font_h":   "Arial", "font_b": "Calibri",
    },
    "codex": {
        "bg":       "F7F0DC", "dark":     "2C1810",
        "primary":  "7B4F1E", "secondary":"C4952A",
        "accent":   "5A8F3C", "muted":    "8B7355",
        "text_dark":"F7F0DC", "text_light":"2C1810",
        "card":     "EDE4C8", "card2":    "E0D5B8",
        "warn":     "8B1A1A", "success":  "3A5A28",
        "font_h":   "Bookman Old Style", "font_b": "Calibri",
    },
    "corporate": {
        "bg":       "F8F9FC", "dark":     "1E3A5F",
        "primary":  "1E3A5F", "secondary":"2E86AB",
        "accent":   "F0A500", "muted":    "6B7280",
        "text_dark":"FFFFFF", "text_light":"1A1A2E",
        "card":     "E8EDF5", "card2":    "EDF7FF",
        "warn":     "D97706", "success":  "059669",
        "font_h":   "Calibri", "font_b": "Calibri",
    },
    "detective": {
        "bg":       "D4C5A9", "dark":     "1A1008",
        "primary":  "8B1A1A", "secondary":"1A3A6A",
        "accent":   "C4952A", "muted":    "5A4A3A",
        "text_dark":"D4C5A9", "text_light":"1A1008",
        "card":     "C0AD90", "card2":    "B8A488",
        "warn":     "8B1A1A", "success":  "2D5A1B",
        "font_h":   "Bookman Old Style", "font_b": "Courier New",
    },
    "cyberpunk": {
        "bg":       "0D0D1A", "dark":     "050508",
        "primary":  "00E5FF", "secondary":"FF00AA",
        "accent":   "FFE600", "muted":    "334466",
        "text_dark":"000000", "text_light":"E0EEFF",
        "card":     "111128", "card2":    "0A1A2E",
        "warn":     "FF3300", "success":  "00FF88",
        "font_h":   "Arial", "font_b": "Courier New",
    },
    "noir": {
        "bg":       "1C1C1C", "dark":     "0A0A0A",
        "primary":  "C9A96E", "secondary":"D4B896",
        "accent":   "E8D5A3", "muted":    "444444",
        "text_dark":"0A0A0A", "text_light":"E8E0D4",
        "card":     "252525", "card2":    "2E2E2E",
        "warn":     "CC3333", "success":  "4A7A3A",
        "font_h":   "Bookman Old Style", "font_b": "Calibri",
    },
    "startup": {
        "bg":       "0F172A", "dark":     "030712",
        "primary":  "6366F1", "secondary":"38BDF8",
        "accent":   "F472B6", "muted":    "475569",
        "text_dark":"FFFFFF", "text_light":"F1F5F9",
        "card":     "1E293B", "card2":    "172033",
        "warn":     "EF4444", "success":  "10B981",
        "font_h":   "Calibri", "font_b": "Calibri",
    },
    "scientific": {
        "bg":       "FFFFFF", "dark":     "111827",
        "primary":  "1D4ED8", "secondary":"7C3AED",
        "accent":   "0891B2", "muted":    "6B7280",
        "text_dark":"FFFFFF", "text_light":"111827",
        "card":     "EEF2FF", "card2":    "F0FAFB",
        "warn":     "DC2626", "success":  "059669",
        "font_h":   "Calibri", "font_b": "Calibri",
    },
    "manuscript": {
        "bg":       "FAF7F2", "dark":     "1C1610",
        "primary":  "5C3D1E", "secondary":"3D2B0E",
        "accent":   "8B6914", "muted":    "8B7355",
        "text_dark":"FAF7F2", "text_light":"1C1610",
        "card":     "EDE8DF", "card2":    "E5DDD0",
        "warn":     "8B1A1A", "success":  "2D5A1B",
        "font_h":   "Cambria", "font_b": "Calibri",
    },
    "newspaper": {
        "bg":       "F5F0E8", "dark":     "111111",
        "primary":  "8B0000", "secondary":"333333",
        "accent":   "C4952A", "muted":    "666666",
        "text_dark":"F5F0E8", "text_light":"111111",
        "card":     "E8E3D8", "card2":    "DEDAD0",
        "warn":     "8B0000", "success":  "2D5A1B",
        "font_h":   "Times New Roman", "font_b": "Times New Roman",
    },
    "minimalist": {
        "bg":       "FFFFFF", "dark":     "0A0A0A",
        "primary":  "0A0A0A", "secondary":"444444",
        "accent":   "0A0A0A", "muted":    "AAAAAA",
        "text_dark":"FFFFFF", "text_light":"0A0A0A",
        "card":     "F5F5F5", "card2":    "EEEEEE",
        "warn":     "CC0000", "success":  "007744",
        "font_h":   "Arial", "font_b": "Arial",
    },
}

# Fallback for unknown themes
_DEFAULT_PAL = PALETTES["academic"]

# Layout variant cycle — ensures consecutive slides differ
LAYOUT_CYCLE = [
    "split_right",    # title left, visual right
    "full_header",    # colored full-width header, content below
    "two_column",     # two equal columns
    "card_grid",      # 2-4 cards in a grid
    "split_left",     # visual left, content right
    "stat_spotlight", # large number/stat with context
    "timeline_flow",  # horizontal timeline
    "quote_full",     # full-bleed quote
]


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
#  DATA CLASSES
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

@dataclass
class Topic:
    title: str
    blocks: List[Block]
    section_num: int = 0
    importance: float = 0.0
    topic_type: str = "general"  # general|entities|timeline|data|warning|quote

    @property
    def entities(self): return [b for b in self.blocks if b.type == BlockType.ENTITY]
    @property
    def warnings(self): return [b for b in self.blocks if b.type == BlockType.WARNING]
    @property
    def callouts(self): return [b for b in self.blocks if b.type == BlockType.CALLOUT]
    @property
    def timelines(self): return [b for b in self.blocks if b.type == BlockType.TIMELINE_EVENT]
    @property
    def lists(self): return [b for b in self.blocks if b.type == BlockType.LIST]
    @property
    def tables(self): return [b for b in self.blocks if b.type == BlockType.TABLE]
    @property
    def paragraphs(self): return [b for b in self.blocks if b.type == BlockType.PARAGRAPH]
    @property
    def quotes(self): return [b for b in self.blocks if b.type == BlockType.QUOTE]
    @property
    def definitions(self): return [b for b in self.blocks if b.type == BlockType.DEFINITION]
    @property
    def relationships(self): return [b for b in self.blocks if b.type == BlockType.RELATIONSHIP]
    @property
    def tasks(self): return [b for b in self.blocks if b.type == BlockType.TASK]


@dataclass
class SlideData:
    slide_type: str
    layout: str = "split_right"
    title: str = ""
    subtitle: str = ""
    bullets: List[str] = field(default_factory=list)
    sub_bullets: Dict[int, List[str]] = field(default_factory=dict)
    body_text: str = ""
    entities: List[Dict] = field(default_factory=list)
    table: List[List[str]] = field(default_factory=list)
    timeline: List[Dict] = field(default_factory=list)
    quote: str = ""
    quote_author: str = ""
    warning: str = ""
    stats: List[Dict] = field(default_factory=list)   # [{val, label, icon}]
    chart_bytes: Optional[bytes] = None               # PNG chart embedded
    chart_type: str = ""
    notes: str = ""
    icon: str = ""
    section_label: str = ""
    is_dark: bool = False       # dark bg slide (title/conclusion)
    visual_hint: str = ""          # passed to visual generator
    visual_bytes: Optional[bytes] = None  # pre-rendered visual PNG


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
#  TOPIC DETECTOR
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

class TopicDetector:
    def detect(self, doc: Document) -> List[Topic]:
        topics: List[Topic] = []
        for i, section in enumerate(doc.sections):
            blocks = self._collect_blocks(section)
            if not blocks:
                continue
            t = Topic(
                title=section.title,
                blocks=blocks,
                section_num=i + 1,
                importance=self._importance(blocks),
                topic_type=self._classify(blocks),
            )
            topics.append(t)

        # Fallback: group by headings if no sections
        if not topics:
            topics = self._group_flat(doc.all_blocks)

        return topics

    def _collect_blocks(self, section: Section) -> List[Block]:
        blocks = [b for b in section.blocks if b.type != BlockType.HEADING]
        for sub in section.subsections:
            blocks += [b for b in sub.blocks if b.type != BlockType.HEADING]
        return blocks

    def _importance(self, blocks: List[Block]) -> float:
        if not blocks:
            return 0.0
        return sum(b.importance_score for b in blocks) / len(blocks)

    def _classify(self, blocks: List[Block]) -> str:
        ents = sum(1 for b in blocks if b.type == BlockType.ENTITY)
        timels = sum(1 for b in blocks if b.type == BlockType.TIMELINE_EVENT)
        warns = sum(1 for b in blocks if b.type == BlockType.WARNING)
        tables = sum(1 for b in blocks if b.type == BlockType.TABLE)
        quotes = sum(1 for b in blocks if b.type == BlockType.QUOTE)
        if ents >= 3: return "entities"
        if timels >= 3: return "timeline"
        if tables >= 1: return "data"
        if warns >= 1: return "warning"
        if quotes >= 1: return "quote"
        return "general"

    def _group_flat(self, blocks: List[Block]) -> List[Topic]:
        groups: List[Tuple[str, List[Block]]] = []
        title, current = "Overview", []
        for b in blocks:
            if b.type == BlockType.HEADING and b.level <= 2:
                if current:
                    groups.append((title, current))
                title, current = b.content, []
            else:
                current.append(b)
        if current:
            groups.append((title, current))
        return [
            Topic(title=t, blocks=bs, section_num=i+1,
                  importance=self._importance(bs), topic_type=self._classify(bs))
            for i, (t, bs) in enumerate(groups) if bs
        ]


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
#  CHART BUILDER  (matplotlib → PNG bytes)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

class ChartBuilder:
    def bar_chart(self, labels: List[str], values: List[float],
                  title: str, pal: Dict, horizontal: bool = False) -> bytes:
        try:
            import matplotlib
            matplotlib.use("Agg")
            import matplotlib.pyplot as plt

            bg     = "#" + pal["bg"]
            prim   = "#" + pal["primary"]
            sec    = "#" + pal["secondary"]
            txtcol = "#" + pal["text_light"]

            n = len(labels)
            colors = [prim if i % 2 == 0 else sec for i in range(n)]

            fig, ax = plt.subplots(figsize=(5.5, 3.2), facecolor=bg)
            ax.set_facecolor(bg)

            if horizontal:
                bars = ax.barh(labels, values, color=colors, height=0.55)
                ax.set_xlabel("", color=txtcol)
                ax.invert_yaxis()
                for bar in bars:
                    w = bar.get_width()
                    ax.text(w * 1.01, bar.get_y() + bar.get_height()/2,
                            f"{w:.0f}", va="center", ha="left",
                            fontsize=8, color=txtcol)
            else:
                bars = ax.bar(labels, values, color=colors, width=0.55)
                for bar in bars:
                    h = bar.get_height()
                    ax.text(bar.get_x() + bar.get_width()/2, h * 1.01,
                            f"{h:.0f}", ha="center", va="bottom",
                            fontsize=8, color=txtcol)

            ax.tick_params(colors=txtcol, labelsize=8)
            ax.spines[:].set_visible(False)
            ax.tick_params(length=0)
            for spine in ax.spines.values():
                spine.set_color(txtcol + "30")
            ax.set_title(title, color=txtcol, fontsize=10, pad=8, fontweight="bold")
            plt.tight_layout(pad=0.5)

            buf = io.BytesIO()
            plt.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                        facecolor=bg, transparent=False)
            plt.close(fig)
            return buf.getvalue()
        except Exception:
            return b""

    def pie_chart(self, labels: List[str], values: List[float],
                  title: str, pal: Dict) -> bytes:
        try:
            import matplotlib
            matplotlib.use("Agg")
            import matplotlib.pyplot as plt

            bg    = "#" + pal["bg"]
            prim  = "#" + pal["primary"]
            sec   = "#" + pal["secondary"]
            acc   = "#" + pal["accent"]
            txtcol= "#" + pal["text_light"]

            # Build color cycle from theme
            base_colors = [prim, sec, acc,
                           "#" + pal["success"], "#" + pal["muted"]]
            colors = (base_colors * math.ceil(len(labels) / len(base_colors)))[:len(labels)]

            fig, ax = plt.subplots(figsize=(4.5, 3.0), facecolor=bg)
            ax.set_facecolor(bg)
            wedges, texts, autotexts = ax.pie(
                values, labels=labels, autopct="%1.0f%%",
                colors=colors, startangle=90,
                textprops={"color": txtcol, "fontsize": 8},
                pctdistance=0.78,
            )
            for at in autotexts:
                at.set_color("#" + pal["text_dark"])
                at.set_fontsize(7)
            ax.set_title(title, color=txtcol, fontsize=10, fontweight="bold", pad=4)
            plt.tight_layout(pad=0.3)

            buf = io.BytesIO()
            plt.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                        facecolor=bg, transparent=False)
            plt.close(fig)
            return buf.getvalue()
        except Exception:
            return b""

    def extract_numeric_data(self, blocks: List[Block]) -> Optional[Tuple[List[str], List[float], str]]:
        """Try to extract chart-ready data from lists or tables."""
        # Try tables first
        for b in blocks:
            if b.type == BlockType.TABLE:
                rows = [r for r in b.content.splitlines()
                        if r.strip() and not re.match(r"^\|[-:\s|]+\|?$", r)]
                if len(rows) >= 3:
                    data_rows = rows[1:]  # skip header
                    labels, vals = [], []
                    for row in data_rows[:8]:
                        cells = [c.strip() for c in row.strip("|").split("|")]
                        if len(cells) >= 2:
                            # Try to parse second cell as number
                            num_str = re.sub(r"[^\d.\-]", "", cells[1])
                            try:
                                vals.append(float(num_str))
                                labels.append(cells[0][:20])
                            except ValueError:
                                pass
                    if len(vals) >= 3:
                        return labels, vals, rows[0].strip("|").split("|")[0].strip()

        # Try lists with numbers
        for b in blocks:
            if b.type == BlockType.LIST:
                labels, vals = [], []
                for line in b.content.splitlines():
                    line = line.strip().lstrip("-•*").strip()
                    # Pattern: "Label: 42" or "Label — 42"
                    m = re.match(r"^(.+?)[\s:—\-]+(\d[\d,.]*)%?\s*$", line)
                    if m:
                        try:
                            v = float(m.group(2).replace(",", ""))
                            labels.append(m.group(1).strip()[:20])
                            vals.append(v)
                        except ValueError:
                            pass
                if len(vals) >= 3:
                    return labels, vals, ""

        return None


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
#  SLIDE PLANNER
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

class SlidePlanner:
    ICONS = {
        # Project types
        "Story Bible": "📖", "Research Notes": "🔬", "Project Plan": "📋",
        "Knowledge Base": "🧠", "Worldbuilding Document": "🌍",
        "Study Notes": "📝", "Document": "📄", "Debate Preparation": "⚖",
        "Competition Planning": "🏆",
        # Topic types
        "entities": "👤", "timeline": "📅", "data": "📊",
        "warning": "⚠️", "quote": "💬", "general": "◆",
        # Content
        "summary": "✅", "agenda": "📋", "conclusion": "🎯",
    }

    def __init__(self):
        self.charts = ChartBuilder()
        self._layout_idx = 0

    def _next_layout(self) -> str:
        l = LAYOUT_CYCLE[self._layout_idx % len(LAYOUT_CYCLE)]
        self._layout_idx += 1
        return l

    def plan(self, doc: Document, topics: List[Topic], mode: str, pal: Dict) -> List[SlideData]:
        self._layout_idx = 0
        slides: List[SlideData] = []

        # ── 1. Title slide (dark) ──────────────────────────────
        subtitle_parts = [doc.project_type.value]
        entity_count = len(doc.get_blocks_by_type(BlockType.ENTITY))
        if entity_count:
            subtitle_parts.append(f"{entity_count} entities")
        warn_count = len(doc.get_blocks_by_type(BlockType.WARNING))
        if warn_count:
            subtitle_parts.append(f"{warn_count} alerts")

        slides.append(SlideData(
            slide_type="title",
            layout="title_dark",
            title=doc.title,
            subtitle=" · ".join(subtitle_parts),
            body_text=self._title_tagline(doc),
            icon=self.ICONS.get(doc.project_type.value, "📄"),
            is_dark=True,
            notes=f"Welcome. This presentation covers {doc.title}. {len(topics)} major sections.",
        ))

        # ── 2. Agenda ──────────────────────────────────────────
        if len(topics) >= 2:
            agenda_bullets = [t.title for t in topics[:10]]
            slides.append(SlideData(
                slide_type="agenda",
                layout="agenda_grid",
                title="What We'll Cover",
                bullets=agenda_bullets,
                notes="Roadmap for this session: " + ", ".join(agenda_bullets[:5]),
            ))

        # ── 3. Content slides ──────────────────────────────────
        for topic in topics:
            slides.extend(self._plan_topic(topic, mode, pal, doc))

        # ── 4. Summary/conclusion (dark) ──────────────────────
        key_points = self._extract_key_points(topics, mode)
        top_entities = sorted(doc.get_blocks_by_type(BlockType.ENTITY),
                              key=lambda b: b.importance_score, reverse=True)[:5]

        slides.append(SlideData(
            slide_type="summary",
            layout="summary_dark",
            title="Key Takeaways",
            bullets=key_points,
            entities=[{"name": e.content, "desc": "", "score": e.importance_score}
                      for e in top_entities],
            icon=self.ICONS["summary"],
            is_dark=True,
            notes="Summary of the main points covered.",
        ))

        return slides

    def _plan_topic(self, topic: Topic, mode: str, pal: Dict, doc: Document) -> List[SlideData]:
        slides: List[SlideData] = []
        label = f"§{topic.section_num}"

        # ── Section title card ─────────────────────────────────
        # (skip if it'll immediately be followed by a richer slide)
        if topic.topic_type in ("general", "quote") or len(topic.blocks) <= 3:
            # Merge into a single richer slide below instead
            pass
        else:
            # Light section-header slide only for big topics
            pass  # We'll handle within the main slide

        # ── Main content slide ─────────────────────────────────
        layout = self._next_layout()
        main = self._build_main_slide(topic, mode, pal, layout, label, doc)
        slides.append(main)

        # ── Supplementary slides ───────────────────────────────

        # Entity showcase (if many entities)
        if topic.topic_type == "entities" and len(topic.entities) >= 3:
            slides.append(self._entity_slide(topic, mode, label))

        # Timeline
        if len(topic.timelines) >= 3:
            slides.append(self._timeline_slide(topic, label))

        # Data / chart
        if topic.tables or (topic.topic_type == "data"):
            chart_slide = self._data_slide(topic, mode, pal, label)
            if chart_slide:
                slides.append(chart_slide)

        # Warning callout
        high_warns = [b for b in topic.warnings if b.importance_score >= 50]
        if high_warns:
            slides.append(self._warning_slide(topic, high_warns, label))

        # Quote slide
        good_quotes = [b for b in topic.quotes if b.importance_score >= 40]
        if good_quotes and len(good_quotes[0].content) > 30:
            slides.append(self._quote_slide(good_quotes[0], label))

        # Detailed evidence slide (detailed mode only)
        if mode == "detailed" and len(topic.paragraphs) >= 2:
            slides.append(self._evidence_slide(topic, label))

        # Stats spotlight if numeric data found
        stats = self._extract_stats(topic)
        if stats and len(stats) >= 2:
            slides.append(SlideData(
                slide_type="stats",
                layout="stat_spotlight",
                title=f"{topic.title} — By the Numbers",
                stats=stats[:4],
                section_label=label,
                notes=f"Key statistics for {topic.title}.",
            ))

        return slides

    def _build_main_slide(self, topic: Topic, mode: str, pal: Dict,
                           layout: str, label: str, doc: Document) -> SlideData:
        """The primary slide for a topic."""
        bullets, sub_bullets = self._extract_bullets(topic, mode)
        icon = self.ICONS.get(topic.topic_type, "◆")

        # Try chart if data available
        chart_bytes = None
        chart_type = ""
        numeric = self.charts.extract_numeric_data(topic.blocks)
        if numeric:
            labels, vals, title = numeric
            if len(labels) <= 6:
                chart_bytes = self.charts.pie_chart(labels, vals, title or topic.title, pal)
                chart_type = "pie"
            else:
                chart_bytes = self.charts.bar_chart(labels, vals, title or topic.title, pal,
                                                    horizontal=len(labels) > 5)
                chart_type = "bar"

        # Choose layout based on whether we have a chart
        if chart_bytes:
            layout = "split_right"
        elif topic.topic_type == "entities" and len(topic.entities) >= 2:
            layout = "card_grid"
        elif topic.topic_type == "timeline":
            layout = "timeline_flow"

        # Build entity cards for card_grid layout
        entities = []
        if layout == "card_grid":
            for e in topic.entities[:4]:
                desc = self._entity_desc(e, topic)
                entities.append({"name": e.content, "desc": desc, "score": e.importance_score})
            # If no entities, fall back
            if not entities:
                layout = "full_header"

        # Visual type selection
        visual_types = ["abstract_geo","icon_cluster","wave_pattern","grid_pattern","data_glyph","topic_art"]
        v_type = visual_types[topic.section_num % len(visual_types)]
        if chart_bytes:
            v_bytes = chart_bytes  # already have a chart, use it
        else:
            v_bytes = gen_visual(v_type, topic.title, pal)

        return SlideData(
            slide_type="topic",
            layout=layout,
            title=topic.title,
            bullets=bullets,
            sub_bullets=sub_bullets,
            entities=entities,
            chart_bytes=chart_bytes,
            chart_type=chart_type,
            icon=icon,
            section_label=label,
            visual_hint=v_type,
            visual_bytes=v_bytes,
            notes=f"Section: {topic.title}. " + (" | ".join(bullets[:3]) if bullets else ""),
        )

    def _entity_slide(self, topic: Topic, mode: str, label: str) -> SlideData:
        entities = []
        for e in topic.entities[:6]:
            entities.append({
                "name": e.content,
                "desc": self._entity_desc(e, topic),
                "score": e.importance_score,
            })
        return SlideData(
            slide_type="entities",
            layout="card_grid",
            title=f"{topic.title} — Key Entities",
            entities=entities,
            icon=self.ICONS["entities"],
            section_label=label,
            notes=f"{len(entities)} entities: " + ", ".join(e["name"] for e in entities),
        )

    def _timeline_slide(self, topic: Topic, label: str) -> SlideData:
        events = []
        for b in topic.timelines[:8]:
            parts = b.content.split(":", 1)
            events.append({
                "year": parts[0].strip() if len(parts) == 2 else "·",
                "event": (parts[1].strip() if len(parts) == 2 else b.content)[:80],
                "score": b.importance_score,
            })
        return SlideData(
            slide_type="timeline",
            layout="timeline_flow",
            title=f"{topic.title} — Timeline",
            timeline=events,
            icon=self.ICONS["timeline"],
            section_label=label,
            notes=f"Timeline with {len(events)} events.",
        )

    def _data_slide(self, topic: Topic, mode: str, pal: Dict, label: str) -> Optional[SlideData]:
        numeric = self.charts.extract_numeric_data(topic.blocks)
        chart_bytes = None
        table = []

        if numeric:
            labels, vals, chart_title = numeric
            if len(labels) <= 5:
                chart_bytes = self.charts.bar_chart(labels, vals, chart_title or topic.title, pal)
            else:
                chart_bytes = self.charts.bar_chart(labels, vals, chart_title or topic.title, pal,
                                                    horizontal=True)

        if topic.tables:
            table = self._parse_table(topic.tables[0].content)

        if not chart_bytes and not table:
            return None

        return SlideData(
            slide_type="data",
            layout="split_right",
            title=f"{topic.title} — Data",
            table=table[:8] if table else [],
            chart_bytes=chart_bytes,
            icon=self.ICONS["data"],
            section_label=label,
            notes=f"Data visualization for {topic.title}.",
        )

    def _warning_slide(self, topic: Topic, warns: List[Block], label: str) -> SlideData:
        return SlideData(
            slide_type="warning",
            layout="warning_full",
            title="⚠️  Alert",
            warning=warns[0].content,
            bullets=[b.content[:140] for b in warns[1:4]],
            icon="⚠️",
            section_label=label,
            notes=f"Warning: {warns[0].content[:80]}",
        )

    def _quote_slide(self, quote_block: Block, label: str) -> SlideData:
        q = quote_block.content
        author = ""
        for sep in [" — ", " - ", " ~ ", "\n— ", "\n- "]:
            if sep in q:
                q, author = q.rsplit(sep, 1)
                break
        return SlideData(
            slide_type="quote",
            layout="quote_full",
            title="",
            quote=q.strip('"\'').strip(),
            quote_author=author.strip(),
            section_label=label,
            notes=f'Quote: "{q[:80]}"',
        )

    def _evidence_slide(self, topic: Topic, label: str) -> SlideData:
        best_paras = sorted(topic.paragraphs, key=lambda b: b.importance_score, reverse=True)[:3]
        return SlideData(
            slide_type="evidence",
            layout="two_column",
            title=f"{topic.title} — In Depth",
            bullets=[p.content[:200] for p in best_paras],
            section_label=label,
            notes=f"Detailed evidence for {topic.title}.",
        )

    # ── Helpers ───────────────────────────────────────────────

    def _extract_bullets(self, topic: Topic, mode: str) -> Tuple[List[str], Dict]:
        bullets: List[str] = []
        sub: Dict[int, List[str]] = {}
        max_b = 3 if mode == "presenter" else 5
        max_len = 65 if mode == "presenter" else 130

        # Callouts first
        for b in topic.callouts[:2]:
            bullets.append(b.content[:max_len])

        # Lists
        for lb in topic.lists[:1]:
            items = [i.lstrip("-•* ").strip() for i in lb.content.splitlines() if i.strip()]
            for it in items[:5]:
                if it:
                    bullets.append(it[:max_len])

        # Tasks (checklists) -- flatten each task + its nested subtasks
        for tb in topic.tasks[:2]:
            for text, checked in self._flatten_task_bullets(
                    tb.content, bool(tb.metadata.get("checked")), tb.metadata.get("subtasks") or []):
                box = "[x] " if checked else "[ ] "
                bullets.append((box + text)[:max_len])
                if len(bullets) >= max_b:
                    break

        # Definitions
        for b in topic.definitions[:2]:
            if ":" in b.content:
                term, defn = b.content.split(":", 1)
                idx = len(bullets)
                bullets.append(term.strip()[:60])
                sub[idx] = [defn.strip()[:120]]
            else:
                bullets.append(b.content[:max_len])

        # Relationships as bullets
        for b in topic.relationships[:2]:
            bullets.append(b.content[:max_len])

        # Paragraphs (detailed or thin content)
        if mode == "detailed" or len(bullets) < 2:
            for p in sorted(topic.paragraphs, key=lambda b: b.importance_score, reverse=True)[:3]:
                bullets.append(p.content[:max_len])

        return bullets[:max_b], sub

    def _flatten_task_bullets(self, text, checked, subtasks):
        out = [(text, checked)]
        for s in subtasks or []:
            out.extend(self._flatten_task_bullets(s.get("text", ""), bool(s.get("checked")),
                                                    s.get("subtasks") or []))
        return out

    def _entity_desc(self, entity: Block, topic: Topic) -> str:
        try:
            idx = topic.blocks.index(entity)
        except ValueError:
            return ""
        for b in topic.blocks[idx+1:idx+4]:
            if b.type in (BlockType.PARAGRAPH, BlockType.RELATIONSHIP, BlockType.DEFINITION):
                return b.content[:120]
        return ""

    def _extract_stats(self, topic: Topic) -> List[Dict]:
        stats = []
        # Look for numbers in paragraphs and lists only -- TABLE/CODE blocks
        # hold raw markdown/code syntax (not prose), and TIMELINE_EVENT
        # dates already get their own dedicated timeline slide, so scanning
        # any of those here just dumps garbled syntax or duplicate dates
        # onto this slide instead of genuine callout statistics.
        prose_types = (BlockType.PARAGRAPH, BlockType.LIST, BlockType.CALLOUT)
        for b in topic.blocks:
            if b.type not in prose_types:
                continue
            for m in re.finditer(r"\b(\d[\d,]*(?:\.\d+)?)\s*(%|x|X|\+)?\b\s*([A-Za-z][^.\n]{3,40})?", b.content):
                val = m.group(1)
                unit = m.group(2) or ""
                label = (m.group(3) or "").strip()[:30] or b.content[:30]
                if float(val.replace(",", "")) > 0:
                    stats.append({"val": val + unit, "label": label, "icon": "◆"})
                    if len(stats) >= 6:
                        break
            if len(stats) >= 6:
                break
        # Deduplicate
        seen = set()
        result = []
        for s in stats:
            k = s["val"]
            if k not in seen:
                seen.add(k)
                result.append(s)
        return result[:4]

    def _extract_key_points(self, topics: List[Topic], mode: str) -> List[str]:
        points = []
        for t in topics[:7]:
            # Best sentence from top callout or paragraph
            best = sorted(
                t.callouts + t.paragraphs,
                key=lambda b: b.importance_score, reverse=True
            )
            if best:
                text = best[0].content[:100]
                points.append(f"{t.title}: {text}")
            else:
                points.append(t.title)
        return points[:6]

    def _parse_table(self, raw: str) -> List[List[str]]:
        rows = []
        for line in raw.strip().splitlines():
            line = line.strip()
            if not line or re.match(r"^\|[-\s|:]+\|$", line):
                continue
            cells = [c.strip() for c in line.strip("|").split("|")]
            if cells:
                rows.append(cells)
        return rows

    def _title_tagline(self, doc: Document) -> str:
        high = sorted(
            [b for b in doc.all_blocks if b.type in (BlockType.PARAGRAPH, BlockType.CALLOUT)],
            key=lambda b: b.importance_score, reverse=True,
        )
        if high:
            return high[0].content[:160]
        return ""


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
#  PPTX RENDERER
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

class PptxRenderer:
    W = 13.33   # inches
    H = 7.5

    def render(self, slides: List[SlideData], pal: Dict, mode: str,
               ornament_text: str = None) -> bytes:
        from pptx import Presentation as Pptx
        from pptx.util import Inches

        prs = Pptx()
        prs.slide_width  = Inches(self.W)
        prs.slide_height = Inches(self.H)
        blank = prs.slide_layouts[6]

        for i, sd in enumerate(slides):
            ps = prs.slides.add_slide(blank)

            # Background
            bg_color = pal["dark"] if sd.is_dark else pal["bg"]
            self._bg(ps, bg_color)

            # Speaker notes
            if sd.notes:
                ps.notes_slide.notes_text_frame.text = sd.notes

            # Route to renderer
            fn = getattr(self, f"_render_{sd.slide_type}", self._render_topic)
            fn(ps, sd, pal, mode, i)

            # Decoration: the same top/bottom-margin ornament motif used by
            # the HTML/PDF/DOCX exporters (decorations.margin_text), placed
            # as a small bottom-right corner mark. Skipped on the title
            # slide (which already carries its own "Information Architecture
            # Studio" footer branding in that corner) and on slide types
            # that already draw a full-width accent bar flush with the
            # bottom edge (quote/summary), so the ornament never overlaps
            # existing slide furniture.
            if ornament_text and sd.slide_type not in ("title", "quote", "summary"):
                txt_color = pal["text_dark"] if sd.is_dark else pal["text_light"]
                self._text(ps, self.W - 3.2, self.H - 0.34, 3.0, 0.28,
                           ornament_text, size=8, color=txt_color + "99",
                           align="right", font=pal["font_b"])

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    # ─── Slide type renderers ─────────────────────────────────

    def _render_title(self, ps, sd: SlideData, pal, mode, idx):
        P    = pal["primary"]
        txt  = pal["text_dark"]

        # Left text panel
        self._rect(ps, 0, 0, self.W * 0.52, self.H, P)
        # Right visual panel (darker shade)
        self._rect(ps, self.W * 0.52, 0, self.W * 0.48, self.H, pal["dark"])
        # Thin accent divider
        self._rect(ps, self.W * 0.52 - 0.04, 0, 0.08, self.H, pal["accent"])

        # Visual image in right panel
        if sd.visual_bytes:
            self._image(ps, self.W * 0.52 + 0.15, 0.25, self.W * 0.46 - 0.3, self.H - 0.5, sd.visual_bytes)
        elif sd.icon:
            self._text(ps, self.W * 0.59, self.H * 0.25, self.W * 0.35, self.H * 0.5,
                       sd.icon, size=80, align="center", color=pal["text_dark"])

        # Title text
        size = 48 if len(sd.title) < 25 else (38 if len(sd.title) < 40 else 30)
        self._text(ps, 0.6, 1.2, self.W * 0.48, 3.2, sd.title,
                   size=size, bold=True, color=txt, wrap=True, font=pal["font_h"])

        # Subtitle pill
        if sd.subtitle:
            self._rect(ps, 0.6, 4.65, min(5.0, len(sd.subtitle) * 0.12 + 0.8), 0.44, pal["accent"])
            self._text(ps, 0.72, 4.69, 4.8, 0.36, sd.subtitle.upper(),
                       size=11, bold=True, color=pal["text_dark"], font=pal["font_b"])

        # Tagline
        if sd.body_text:
            self._text(ps, 0.6, 5.3, self.W * 0.47, 1.0, sd.body_text[:100],
                       size=13, color=txt + "BB", wrap=True, font=pal["font_b"])

        self._text(ps, 0.6, self.H - 0.42, 4.0, 0.3,
                   "Information Architecture Studio",
                   size=8, color=txt + "66", font=pal["font_b"])

    def _render_agenda(self, ps, sd: SlideData, pal, mode, idx):
        P, txt = pal["primary"], pal["text_light"]
        items = sd.bullets

        # Left column: large title
        self._rect(ps, 0, 0, 4.2, self.H, P)
        self._text(ps, 0.4, 1.5, 3.4, 3.0, sd.title,
                   size=36, bold=True, color=pal["text_dark"], wrap=True, font=pal["font_h"])
        self._text(ps, 0.4, self.H - 0.55, 3.4, 0.4,
                   f"{len(items)} sections", size=12, color=pal["text_dark"] + "AA", font=pal["font_b"])

        # Right column: numbered items
        n = len(items)
        row_h = min(0.72, (self.H - 1.0) / max(n, 1))
        for j, item in enumerate(items):
            y = 0.5 + j * row_h
            # Number circle
            self._circle(ps, 4.6, y + 0.08, 0.32, P)
            self._text(ps, 4.58, y + 0.06, 0.36, 0.32, str(j + 1),
                       size=12, bold=True, color=pal["text_dark"], align="center", font=pal["font_b"])
            self._text(ps, 5.12, y, 7.7, row_h, item,
                       size=15 if n <= 6 else 13, color=txt, wrap=True, font=pal["font_b"])

    def _render_topic(self, ps, sd: SlideData, pal, mode, idx):
        layout = sd.layout
        if layout == "card_grid" and sd.entities:
            self._layout_card_grid(ps, sd, pal, mode)
        elif layout == "split_right" and sd.chart_bytes:
            self._layout_chart_split(ps, sd, pal, mode)
        elif layout == "two_column" and sd.bullets and len(sd.bullets) >= 4:
            self._layout_two_column(ps, sd, pal, mode)
        elif layout == "split_left":
            self._layout_split_left(ps, sd, pal, mode)
        elif layout == "stat_spotlight" and sd.stats:
            self._layout_stats(ps, sd, pal, mode)
        elif layout == "timeline_flow" and sd.timeline:
            self._layout_timeline(ps, sd, pal, mode)
        else:
            self._layout_full_header(ps, sd, pal, mode)

    def _render_entities(self, ps, sd: SlideData, pal, mode, idx):
        self._layout_card_grid(ps, sd, pal, mode)

    def _render_timeline(self, ps, sd: SlideData, pal, mode, idx):
        self._layout_timeline(ps, sd, pal, mode)

    def _render_data(self, ps, sd: SlideData, pal, mode, idx):
        if sd.chart_bytes:
            self._layout_chart_split(ps, sd, pal, mode)
        elif sd.table:
            self._layout_table(ps, sd, pal, mode)
        else:
            self._layout_full_header(ps, sd, pal, mode)

    def _render_warning(self, ps, sd: SlideData, pal, mode, idx):
        WARN = pal["warn"]
        WARNL = self._lighten(WARN, 0.88)

        self._rect(ps, 0, 0, self.W, self.H, WARNL)
        # Bold left stripe — intentionally not a "decorative bar", it's a functional severity indicator
        self._rect(ps, 0, 0, 0.18, self.H, WARN)
        # Alert header box
        self._rect(ps, 0.35, 0.3, self.W - 0.5, 1.4, WARN)
        self._text(ps, 0.5, 0.42, self.W - 0.7, 1.1, sd.title,
                   size=34, bold=True, color="FFFFFF", font=pal["font_h"])

        if sd.warning:
            self._rect(ps, 0.35, 1.9, self.W - 0.5, 2.2, "FFFFFF" + "BB")
            self._text(ps, 0.55, 2.0, self.W - 0.8, 2.0, sd.warning,
                       size=18, color=self._darken(WARN, 0.3), wrap=True, font=pal["font_b"])

        for j, b in enumerate(sd.bullets[:4]):
            y = 4.3 + j * 0.7
            self._circle(ps, 0.45, y + 0.1, 0.22, WARN)
            self._text(ps, 0.82, y, self.W - 1.1, 0.65,
                       b, size=14, color=pal["text_light"], wrap=True, font=pal["font_b"])

    def _render_quote(self, ps, sd: SlideData, pal, mode, idx):
        P = pal["primary"]
        BG = pal["dark"]
        txt = pal["text_dark"]

        self._rect(ps, 0, 0, self.W, self.H, BG)
        self._rect(ps, 0, 0, self.W, 0.06, P)
        self._rect(ps, 0, self.H - 0.06, self.W, 0.06, P)

        # Decorative quote marks
        self._text(ps, 0.4, 0.4, 2.0, 2.2, "\u201c",
                   size=120, color=P + "55", bold=True, font=pal["font_h"])

        q = sd.quote
        font_size = 30 if len(q) < 80 else (24 if len(q) < 150 else 19)
        self._text(ps, 1.0, 1.3, self.W - 1.8, 4.2, q,
                   size=font_size, color=txt, wrap=True, align="center", font=pal["font_h"])

        self._text(ps, 8.0, 5.1, 4.8, 1.2, "\u201d",
                   size=90, color=P + "55", bold=True, font=pal["font_h"])

        if sd.quote_author:
            self._text(ps, 1.0, self.H - 1.0, self.W - 2.0, 0.6,
                       f"— {sd.quote_author}", size=15,
                       color=P, align="center", bold=True, font=pal["font_b"])

        if sd.section_label:
            self._text(ps, self.W - 1.4, 0.12, 1.2, 0.28, sd.section_label,
                       size=9, color=txt + "66", align="right", font=pal["font_b"])

    def _render_evidence(self, ps, sd: SlideData, pal, mode, idx):
        self._layout_two_column(ps, sd, pal, mode)

    def _render_stats(self, ps, sd: SlideData, pal, mode, idx):
        self._layout_stats(ps, sd, pal, mode)

    def _render_summary(self, ps, sd: SlideData, pal, mode, idx):
        P = pal["primary"]
        S = pal["secondary"]
        txt = pal["text_dark"]

        # Dark background
        self._rect(ps, 0, 0, self.W, self.H, pal["dark"])
        # Accent top band
        self._rect(ps, 0, 0, self.W, 0.08, P)
        self._rect(ps, 0, self.H - 0.08, self.W, 0.08, S)

        self._text(ps, 0.6, 0.25, 6.0, 0.9, sd.title,
                   size=36, bold=True, color=txt, font=pal["font_h"])
        if sd.icon:
            self._text(ps, self.W - 1.4, 0.2, 1.0, 0.9, sd.icon, size=30, color=P)

        # Bullets (left column)
        for j, b in enumerate(sd.bullets[:6]):
            y = 1.4 + j * 0.83
            self._circle(ps, 0.55, y + 0.15, 0.25, P)
            self._text(ps, 0.55, y + 0.1, 0.25, 0.25, str(j+1),
                       size=10, bold=True, color=txt, align="center", font=pal["font_b"])
            self._text(ps, 0.97, y, 6.4, 0.8, b,
                       size=14, color=txt, wrap=True, font=pal["font_b"])

        # Entity chips (right column)
        if sd.entities:
            self._text(ps, 8.1, 1.2, 4.8, 0.38, "KEY ENTITIES",
                       size=10, bold=True, color=P, font=pal["font_b"])
            for j, e in enumerate(sd.entities[:5]):
                y = 1.72 + j * 0.98
                self._rect(ps, 8.1, y, 4.8, 0.82, pal["card"])
                self._text(ps, 8.28, y + 0.15, 4.4, 0.58, e["name"],
                           size=14, bold=True, color=pal["text_light"], font=pal["font_b"])

    # ─── Layout helpers ───────────────────────────────────────

    def _layout_full_header(self, ps, sd: SlideData, pal, mode):
        """Header + bullets left, visual right."""
        P   = pal["primary"]
        txt = pal["text_light"]
        VX  = 7.6   # visual panel x

        # Header band (left 58% only)
        self._rect(ps, 0, 0, VX, 1.4, P)
        self._text(ps, 0.5, 0.18, VX - 0.6, 1.1, sd.title,
                   size=30, bold=True, color=pal["text_dark"], wrap=True, font=pal["font_h"])
        if sd.section_label:
            self._text(ps, VX - 1.4, 0.06, 1.2, 0.26, sd.section_label,
                       size=9, color=pal["text_dark"] + "AA", align="right", font=pal["font_b"])

        # Visual panel right
        self._rect(ps, VX, 0, self.W - VX, self.H, pal["card"])
        if sd.visual_bytes:
            self._image(ps, VX + 0.1, 0.1, self.W - VX - 0.2, self.H - 0.2, sd.visual_bytes)
        elif sd.icon:
            self._text(ps, VX + 0.1, self.H * 0.2, self.W - VX - 0.2, self.H * 0.6,
                       sd.icon, size=52, align="center", color=P)

        # Bullets left panel
        bullets = sd.bullets
        if not bullets:
            self._text(ps, 0.5, 2.2, VX - 0.7, 1.0,
                       "No content available.", size=13, color=pal["muted"], font=pal["font_b"])
            return

        n     = len(bullets)
        row_h = min(1.1, (self.H - 1.7) / max(n, 1))
        for j, b in enumerate(bullets):
            y = 1.6 + j * row_h
            self._circle(ps, 0.44, y + row_h * 0.32, 0.16, P)
            size = 19 if mode == "presenter" else 15
            self._text(ps, 0.76, y, VX - 1.0, row_h,
                       b, size=size, color=txt, wrap=True, font=pal["font_b"])
            if j in sd.sub_bullets:
                sub_y = y + row_h * 0.6
                for sub in sd.sub_bullets[j][:1]:
                    self._text(ps, 1.0, sub_y, VX - 1.3, 0.38,
                               f"  › {sub}", size=11, color=pal["muted"], wrap=True, font=pal["font_b"])
                    sub_y += 0.36

    def _layout_two_column(self, ps, sd: SlideData, pal, mode):
        """Title top, two equal bullet columns."""
        P = pal["primary"]
        txt = pal["text_light"]

        self._rect(ps, 0, 0, self.W, 0.08, P)
        self._text(ps, 0.5, 0.2, self.W - 1.0, 0.9, sd.title,
                   size=30, bold=True, color=txt, font=pal["font_h"])
        if sd.section_label:
            self._text(ps, self.W - 1.5, 0.22, 1.3, 0.3, sd.section_label,
                       size=9, color=pal["muted"], align="right", font=pal["font_b"])

        bullets = sd.bullets
        if not bullets:
            return

        mid = math.ceil(len(bullets) / 2)
        left_b  = bullets[:mid]
        right_b = bullets[mid:]
        row_h = min(0.85, (self.H - 1.4) / max(len(left_b), 1))

        for col, col_bullets in enumerate([left_b, right_b]):
            x_base = 0.4 if col == 0 else 6.8
            for j, b in enumerate(col_bullets):
                y = 1.35 + j * row_h
                self._circle(ps, x_base, y + 0.24, 0.2, P)
                size = 17 if mode == "presenter" else 14
                self._text(ps, x_base + 0.35, y, 5.9, row_h,
                           b, size=size, color=txt, wrap=True, font=pal["font_b"])

        # Visual strip at bottom if space allows and visual available
        if sd.visual_bytes and len(left_b) <= 3:
            self._image(ps, 4.0, 4.5, 5.2, 2.7, sd.visual_bytes)

    def _layout_split_left(self, ps, sd: SlideData, pal, mode):
        """Visual left panel, bullets right."""
        P        = pal["primary"]
        txt      = pal["text_light"]
        VW       = 5.2   # visual panel width

        # Visual panel left
        self._rect(ps, 0, 0, VW, self.H, pal["dark"])
        if sd.visual_bytes:
            self._image(ps, 0.1, 0.1, VW - 0.2, self.H - 0.2, sd.visual_bytes)
        elif sd.icon:
            self._text(ps, 0.1, self.H * 0.25, VW - 0.2, self.H * 0.5,
                       sd.icon, size=72, align="center", color=pal["text_dark"])

        # Content right
        self._text(ps, VW + 0.4, 0.5, self.W - VW - 0.7, 1.2, sd.title,
                   size=26, bold=True, color=txt, wrap=True, font=pal["font_h"])
        bullets = sd.bullets
        row_h   = min(1.05, (self.H - 2.0) / max(len(bullets), 1))
        for j, b in enumerate(bullets):
            y = 2.0 + j * row_h
            self._circle(ps, VW + 0.38, y + 0.28, 0.18, P)
            size = 18 if mode == "presenter" else 14
            self._text(ps, VW + 0.72, y, self.W - VW - 1.0, row_h, b,
                       size=size, color=txt, wrap=True, font=pal["font_b"])

    def _layout_chart_split(self, ps, sd: SlideData, pal, mode):
        """Title left, chart right — or vice versa."""
        P  = pal["primary"]
        txt = pal["text_light"]

        self._text(ps, 0.5, 0.25, 6.8, 1.1, sd.title,
                   size=28, bold=True, color=txt, font=pal["font_h"], wrap=True)
        if sd.section_label:
            self._text(ps, self.W - 1.5, 0.28, 1.3, 0.3, sd.section_label,
                       size=9, color=pal["muted"], align="right", font=pal["font_b"])

        # Chart on right
        if sd.chart_bytes:
            self._image(ps, 6.8, 0.3, 6.2, 6.9, sd.chart_bytes)

        # Bullets on left
        bullets = sd.bullets
        row_h = min(0.88, (self.H - 1.6) / max(len(bullets), 1)) if bullets else 0.88
        for j, b in enumerate(bullets[:5]):
            y = 1.55 + j * row_h
            self._circle(ps, 0.42, y + 0.24, 0.2, P)
            size = 16 if mode == "presenter" else 13
            self._text(ps, 0.77, y, 5.7, row_h, b, size=size,
                       color=txt, wrap=True, font=pal["font_b"])

    def _layout_card_grid(self, ps, sd: SlideData, pal, mode):
        """Grid of entity/concept cards."""
        P   = pal["primary"]
        S   = pal["secondary"]
        txt = pal["text_light"]
        card_col = pal["card"]

        self._text(ps, 0.5, 0.18, self.W - 0.8, 0.9, sd.title,
                   size=28, bold=True, color=txt, font=pal["font_h"])
        if sd.section_label:
            self._text(ps, self.W - 1.5, 0.2, 1.3, 0.3, sd.section_label,
                       size=9, color=pal["muted"], align="right", font=pal["font_b"])
        self._rect(ps, 0.5, 1.12, self.W - 0.8, 0.04, P)

        entities = sd.entities
        n = len(entities)
        if n == 0:
            return

        cols = min(3, n)
        rows = math.ceil(n / cols)
        gutter = 0.18
        cw = (self.W - 0.8 - gutter * (cols - 1)) / cols
        ch = (self.H - 1.4 - gutter * (rows - 1)) / rows

        accent_cycle = [P, S, pal["accent"], pal["success"]]

        for j, ent in enumerate(entities):
            c = j % cols
            r = j // cols
            x = 0.4 + c * (cw + gutter)
            y = 1.3  + r * (ch + gutter)
            acc = accent_cycle[j % len(accent_cycle)]

            # Card background
            self._rect(ps, x, y, cw, ch, card_col)
            # Top color patch
            self._rect(ps, x, y, cw, 0.22, acc)

            # Score bar at bottom
            score_w = cw * min(ent.get("score", 50) / 100, 1.0)
            self._rect(ps, x, y + ch - 0.08, score_w, 0.08, acc + "88")

            # Name
            self._text(ps, x + 0.1, y + 0.28, cw - 0.18, 0.5,
                       ent["name"], size=14, bold=True, color=pal["text_light"], font=pal["font_b"])
            # Desc
            if ent.get("desc") and ch > 1.2:
                self._text(ps, x + 0.1, y + 0.82, cw - 0.18, ch - 1.0,
                           ent["desc"], size=11, color=pal["muted"], wrap=True, font=pal["font_b"])

    def _layout_timeline(self, ps, sd: SlideData, pal, mode):
        """Horizontal timeline with alternating labels."""
        P   = pal["primary"]
        S   = pal["secondary"]
        txt = pal["text_light"]

        self._text(ps, 0.5, 0.2, self.W - 0.8, 0.9, sd.title,
                   size=28, bold=True, color=txt, font=pal["font_h"])
        if sd.section_label:
            self._text(ps, self.W - 1.5, 0.22, 1.3, 0.3, sd.section_label,
                       size=9, color=pal["muted"], align="right", font=pal["font_b"])

        events = sd.timeline
        if not events:
            return

        n  = len(events)
        x0 = 0.7
        x1 = self.W - 0.7
        spine_y = 3.75

        # Spine line
        self._rect(ps, x0, spine_y - 0.015, x1 - x0, 0.03, pal["muted"])

        step = (x1 - x0) / max(n - 1, 1)
        accent_cycle = [P, S, pal["accent"]]

        for j, ev in enumerate(events):
            x = x0 + j * step
            acc = accent_cycle[j % 3]

            # Dot on spine
            self._circle(ps, x - 0.15, spine_y - 0.15, 0.3, acc)

            # Year badge (alternating above/below)
            if j % 2 == 0:
                by = spine_y - 1.55
                ey = spine_y + 0.38
            else:
                by = spine_y + 0.78
                ey = spine_y - 1.18

            # Year box
            yr = ev.get("year", "")
            if yr:
                bw = max(0.9, len(yr) * 0.09 + 0.3)
                self._rect(ps, x - bw/2, by, bw, 0.35, acc)
                self._text(ps, x - bw/2, by + 0.02, bw, 0.31, yr,
                           size=11, bold=True, color=pal["text_dark"],
                           align="center", font=pal["font_b"])

            # Event text
            tw = min(1.6, step * 0.9)
            self._text(ps, x - tw/2, ey, tw, 1.2,
                       ev.get("event", ""), size=10, color=txt,
                       align="center", wrap=True, font=pal["font_b"])

    def _layout_table(self, ps, sd: SlideData, pal, mode):
        """Title + table."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        P   = pal["primary"]
        txt = pal["text_light"]

        self._text(ps, 0.5, 0.2, self.W - 0.8, 0.9, sd.title,
                   size=28, bold=True, color=txt, font=pal["font_h"])

        data = sd.table
        if not data:
            return

        rows = len(data)
        cols = max(len(r) for r in data)
        data = [r + [""] * (cols - len(r)) for r in data]

        shape = ps.shapes.add_table(rows, cols,
                                    Inches(0.5), Inches(1.2),
                                    Inches(self.W - 0.8), Inches(self.H - 1.5))
        tbl = shape.table

        acc_rgb  = self._hrx(P)
        text_rgb = self._hrx(txt)
        card_rgb = self._hrx(pal["card"])

        for ri, row in enumerate(data):
            for ci, cell_text in enumerate(row):
                cell = tbl.cell(ri, ci)
                cell.text = str(cell_text)
                tf = cell.text_frame
                p  = tf.paragraphs[0]
                p.font.size = Pt(12)
                if ri == 0:
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(*acc_rgb)
                else:
                    p.font.color.rgb = RGBColor(*text_rgb)
                    cell.fill.solid()
                    bg = card_rgb if ri % 2 == 0 else self._hrx(pal["bg"])
                    cell.fill.fore_color.rgb = RGBColor(*bg)

    def _layout_stats(self, ps, sd: SlideData, pal, mode):
        """Large stat spotlight cards."""
        P   = pal["primary"]
        S   = pal["secondary"]
        txt = pal["text_light"]
        card= pal["card"]

        self._rect(ps, 0, 0, self.W, 1.2, P)
        self._text(ps, 0.5, 0.2, self.W - 0.8, 0.85, sd.title,
                   size=30, bold=True, color=pal["text_dark"], font=pal["font_h"])

        stats = sd.stats[:4]
        n = len(stats)
        if n == 0:
            return

        cw = (self.W - 0.6 - 0.2 * (n - 1)) / n
        accent_cycle = [P, S, pal["accent"], pal["success"]]

        for j, s in enumerate(stats):
            x = 0.3 + j * (cw + 0.2)
            acc = accent_cycle[j % len(accent_cycle)]
            self._rect(ps, x, 1.4, cw, 5.6, card)
            self._rect(ps, x, 1.4, cw, 0.12, acc)

            # Big number
            vsize = 56 if len(s["val"]) <= 5 else (42 if len(s["val"]) <= 8 else 32)
            self._text(ps, x + 0.1, 1.9, cw - 0.2, 2.0,
                       s["val"], size=vsize, bold=True, color=acc,
                       align="center", font=pal["font_h"])
            # Label
            self._text(ps, x + 0.1, 4.1, cw - 0.2, 1.8,
                       s["label"], size=13, color=txt,
                       align="center", wrap=True, font=pal["font_b"])

    # ─── Drawing primitives ───────────────────────────────────

    def _bg(self, ps, hex_color: str):
        from pptx.dml.color import RGBColor
        fill = ps.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*self._hrx(hex_color))

    def _rect(self, ps, x, y, w, h, hex_color: str):
        from pptx.util import Inches
        from pptx.dml.color import RGBColor
        sh = ps.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(*self._hrx(hex_color))
        sh.line.fill.background()
        return sh

    def _circle(self, ps, x, y, d, hex_color: str):
        """Draw a filled circle (ellipse shape)."""
        from pptx.util import Inches
        from pptx.dml.color import RGBColor
        sh = ps.shapes.add_shape(9, Inches(x), Inches(y), Inches(d), Inches(d))  # 9 = oval
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(*self._hrx(hex_color))
        sh.line.fill.background()
        return sh

    def _text(self, ps, x, y, w, h, text: str, size=14, bold=False,
              color="111111", align="left", wrap=True, font="Calibri"):
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        tb = ps.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        p  = tf.paragraphs[0]
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        try:
            # Strip trailing alpha if 8-char hex
            col = color[:6] if len(color) >= 6 else color
            run.font.color.rgb = RGBColor(*self._hrx(col))
        except Exception:
            pass
        return tb

    def _image(self, ps, x, y, w, h, img_bytes: bytes):
        from pptx.util import Inches
        buf = io.BytesIO(img_bytes)
        ps.shapes.add_picture(buf, Inches(x), Inches(y), Inches(w), Inches(h))

    def _hrx(self, hex_color: str):
        """Hex string → (r, g, b) tuple."""
        h = hex_color.lstrip("#").upper()
        if len(h) == 3:
            h = h[0]*2 + h[1]*2 + h[2]*2
        h = (h + "000000")[:6]
        try:
            return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
        except Exception:
            return (0, 0, 0)

    def _lighten(self, hex_c: str, amt: float) -> str:
        r, g, b = self._hrx(hex_c)
        r2 = min(255, int(r + (255-r) * amt))
        g2 = min(255, int(g + (255-g) * amt))
        b2 = min(255, int(b + (255-b) * amt))
        return f"{r2:02X}{g2:02X}{b2:02X}"

    def _darken(self, hex_c: str, amt: float) -> str:
        r, g, b = self._hrx(hex_c)
        r2 = max(0, int(r * (1 - amt)))
        g2 = max(0, int(g * (1 - amt)))
        b2 = max(0, int(b * (1 - amt)))
        return f"{r2:02X}{g2:02X}{b2:02X}"


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
#  PUBLIC API
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def generate_pptx(doc: Document, theme: str = "academic",
                  mode: str = "presenter", custom_palette: dict = None,
                  decorate: bool = False, decoration_style: str = "auto") -> bytes:
    """Full pipeline → PPTX bytes."""
    pal     = custom_palette or PALETTES.get(theme, _DEFAULT_PAL)
    detector = TopicDetector()
    planner  = SlidePlanner()
    renderer = PptxRenderer()

    topics = detector.detect(doc)
    slides = planner.plan(doc, topics, mode=mode, pal=pal)

    ornament_text = None
    if decorate:
        from app.renderer import decorations
        resolved_style = decorations.resolve_style(decoration_style, doc, theme)
        primary, _secondary = decorations.DECORATION_SETS.get(
            resolved_style, decorations.DECORATION_SETS[decorations.DEFAULT_SET])
        ornament_text = decorations.margin_text(primary)

    return renderer.render(slides, pal=pal, mode=mode, ornament_text=ornament_text)
